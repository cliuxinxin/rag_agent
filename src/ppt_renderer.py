# src/ppt_renderer.py
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt

def generate_ppt_binary(slides_data: list) -> bytes:
    """
    根据结构化数据生成 PPT 二进制流 (内存中)
    返回: bytes 对象
    """
    prs = Presentation() # 初始化空白 PPT

    for slide_info in slides_data:
        s_type = slide_info.get("type", "content")
        
        # 1. 创建页面 (逻辑不变)
        if s_type == "cover":
            slide_layout = prs.slide_layouts[0] 
        elif s_type == "section":
            slide_layout = prs.slide_layouts[2] 
        else:
            slide_layout = prs.slide_layouts[1] 

        slide = prs.slides.add_slide(slide_layout)

        # 2. 填充标题 (逻辑不变)
        title = slide.shapes.title
        if title:
            title.text = slide_info.get("title", "")

        # 3. 填充内容 (逻辑不变)
        if s_type == "cover":
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = slide_info.get("subtitle", "")
        elif s_type == "content":
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                if body_shape.has_text_frame:
                    tf = body_shape.text_frame
                    tf.word_wrap = True
                    bullets = slide_info.get("bullets", [])
                    if isinstance(bullets, list):
                        tf.clear() 
                        for point in bullets:
                            p = tf.add_paragraph()
                            p.text = str(point)
                            p.level = 0
                            p.space_after = Pt(10)

        # 4. 填充备注 (逻辑不变)
        notes_txt = slide_info.get("speaker_notes", "")
        if notes_txt:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes_txt

    # === 核心修改：使用 BytesIO ===
    output = BytesIO()
    prs.save(output)  # 保存到内存流
    output.seek(0)    # 指针回滚到开头
    return output.getvalue() # 返回二进制数据 (bytes)