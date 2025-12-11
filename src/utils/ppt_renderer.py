# src/utils/ppt_renderer.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt

def generate_ppt_file(slides_data: list, filename="output.pptx") -> str:
    """
    根据结构化数据生成 PPTX 文件 (使用默认母版)
    Layout 索引参考 (标准空白模板):
    0: Title Slide (Title + Subtitle)
    1: Title and Content (Title + Bullet points)
    2: Section Header (Title + Text)
    """
    prs = Presentation() # 初始化空白 PPT

    for slide_info in slides_data:
        s_type = slide_info.get("type", "content")
        
        # 1. 创建页面 (根据类型选择布局)
        if s_type == "cover":
            slide_layout = prs.slide_layouts[0] # Title Slide
        elif s_type == "section":
            slide_layout = prs.slide_layouts[2] # Section Header
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content (Default)

        slide = prs.slides.add_slide(slide_layout)

        # 2. 填充标题
        title = slide.shapes.title
        if title:
            title.text = slide_info.get("title", "")

        # 3. 填充内容
        if s_type == "cover":
            # 封面副标题
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = slide_info.get("subtitle", "")
                
        elif s_type == "content":
            # 正文列表
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.word_wrap = True
                
                bullets = slide_info.get("bullets", [])
                if isinstance(bullets, list):
                    # 清空默认的空段落
                    tf.clear() 
                    for point in bullets:
                        p = tf.add_paragraph()
                        p.text = str(point)
                        p.level = 0
                        # 稍微调整一下间距
                        p.space_after = Pt(10)

        # 4. 填充演讲者备注 (Speaker Notes)
        notes_txt = slide_info.get("speaker_notes", "")
        if notes_txt:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes_txt

    # 5. 保存文件
    output_dir = "storage/generated_ppts"
    os.makedirs(output_dir, exist_ok=True)
    full_path = os.path.join(output_dir, filename)
    
    # 防止文件名冲突
    if os.path.exists(full_path):
        base, ext = os.path.splitext(filename)
        import time
        full_path = os.path.join(output_dir, f"{base}_{int(time.time())}{ext}")
        
    prs.save(full_path)
    return full_path